import os

from pptx import Presentation
from docx import Document
from docx.shared import Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH

RECOGNIZED_EXTENSIONS = ['png','jpg','bmp','jpeg','gif']

def pptToDocx(INPUT_FILE,OUTPUT_DIR,MARGIN = 2,IMAGE_WIDTH=18):
    # Save docx with same filename
    OUTPUT_FILE = OUTPUT_DIR+(INPUT_FILE.split('/')[-1].split('.')[0]+'.docx')

    # Load presentation and create document file
    prs = Presentation(INPUT_FILE)
    document = Document()   

    # Set page margins
    for section in document.sections:
        section.top_margin = Cm(MARGIN)
        section.bottom_margin = Cm(MARGIN)
        section.left_margin = Cm(MARGIN)
        section.right_margin = Cm(MARGIN)

    # Image counter
    counter = 0

    isFirstSlide = True
    for slide in prs.slides:
        isFirst = True
        for shape in slide.shapes:
            # Autosave images
            if not shape.has_text_frame:
                # If no image then skip
                try:
                    shape.image
                except:
                    continue

                # Skip unrecognized formats
                if shape.image.ext not in RECOGNIZED_EXTENSIONS:
                    continue

                # Save the image file temporarily to disk
                filename = str(counter)+'.'+shape.image.ext
                filePath = OUTPUT_DIR+'imgs/'+filename

                # Create imgs directory if it does not exist
                if not os.path.exists(OUTPUT_DIR+'imgs'):
                    os.makedirs(OUTPUT_DIR+'imgs')
                
                #print(filePath)
                with open(filePath,'wb') as f:
                    f.write(shape.image.blob)

                # If IMAGE_WIDTH is -1 then do not resize
                p = document.add_picture(filePath)
                if p.width.cm > IMAGE_WIDTH:
                    p.height = Cm(p.height.cm / p.width.cm * IMAGE_WIDTH)
                    p.width = Cm(IMAGE_WIDTH)
                document.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
                counter += 1
                
                continue
            
            for paragraph in shape.text_frame.paragraphs:
                # Make the first paragraph the header
                if isFirst:
                    # Make first paragraph in the first slide the document title
                    if isFirstSlide:
                        document.add_heading(paragraph.text, 0)
                        isFirstSlide = False
                    else:
                        document.add_heading(paragraph.text, 1)
                    isFirst = False
                else:
                    # Ignore empty paragraphs
                    if len(paragraph.text) > 0:
                        document.add_paragraph(paragraph.text)
                        document.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

    # Save the output
    document.save(OUTPUT_FILE)

def test():
    INPUT_FILE = './input/Review1.pptx'
    OUTPUT_DIR = './output/'

    pptToDocx(INPUT_FILE,OUTPUT_DIR)

