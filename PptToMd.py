from pptx import Presentation
import markdown2

INPUT_FILE = './input/Review1.pptx'
OUTPUT_DIR = './output/'
OUTPUT_FILE = OUTPUT_DIR+'Review1'

prs = Presentation(INPUT_FILE)

text = []

counter = 0
isFirstSlide = True
for slide in prs.slides:
    isFirst = True
    for shape in slide.shapes:
        if not shape.has_text_frame:
            filename = str(counter)+'.'+shape.image.ext
            #print(filename)
            with open(OUTPUT_DIR+filename,'wb') as f:
                f.write(shape.image.blob)
            text.append('!['+filename+']('+filename+')\n\n')
            counter += 1
            continue
        if isFirst:
            if isFirstSlide:
                text.append('# ')
                isFirstSlide = False
            else:
                text.append('## ')
            isFirst = False
        for paragraph in shape.text_frame.paragraphs:
            text.append(('* '*paragraph.level)+paragraph.text)
            text.append('\n\n')
        text.append('\n')

dumpString = ''.join(text)

with open(OUTPUT_FILE+'.md','w') as f:
    f.write(dumpString)

with open(OUTPUT_FILE+'.html','w') as f:
    f.write(markdown2.markdown(dumpString)+'<style>body{margin:5%}</style>')
